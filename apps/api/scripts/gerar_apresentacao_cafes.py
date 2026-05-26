"""Gera apresentação executiva (PPTX) com a análise dos cafés T&F 1T26.

Padrão visual MudAção (verde escuro #0a4a3a). Editável no PowerPoint/Keynote.

Uso:
    cd apps/api
    uv run --with openpyxl --with pandas --with python-pptx python scripts/gerar_apresentacao_cafes.py

Output:
    apps/api/tmp/Apresentacao_TF_Cafes_PEC8.pptx
"""

from __future__ import annotations

import io
import math
import sys
from dataclasses import dataclass
from decimal import Decimal
from pathlib import Path

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", line_buffering=True)

import pandas as pd  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Emu, Inches, Pt  # noqa: E402

# =============================================================================
# Engine constants (espelho de packages/engine/src/engine/core.py)
# =============================================================================
JORNADA_6X1 = Decimal("44")
JORNADA_5X2 = Decimal("40")
RATIO_BASE = JORNADA_6X1 / JORNADA_5X2
FATOR_PERDAS = {"pessimista": Decimal("1.10"), "neutro": Decimal("1.05"), "otimista": Decimal("0.98")}
PESO_GANHO_PROD = {"pessimista": Decimal("0.5"), "neutro": Decimal("1.0"), "otimista": Decimal("1.5")}
GANHO_PROD_DEFAULT = Decimal("0.05")
WFM_ECONOMY_PCT = Decimal("0.05")


def _ratio(c: str) -> Decimal:
    return RATIO_BASE * FATOR_PERDAS[c] * (Decimal("1") - GANHO_PROD_DEFAULT * PESO_GANHO_PROD[c])


EXCEL_PATH = Path(r"C:\Users\felip\Downloads\INFORMAÇÃO 1T26 CAFÉ.xlsx")

# =============================================================================
# Paleta MudAção
# =============================================================================
COR_VERDE_ESCURO = RGBColor(0x06, 0x29, 0x20)  # #062920
COR_VERDE_PRIMARIO = RGBColor(0x0A, 0x4A, 0x3A)  # #0a4a3a
COR_VERDE_MEDIO = RGBColor(0x22, 0x55, 0x3D)  # #22553d
COR_VERDE_CLARO = RGBColor(0xB8, 0xDC, 0xC8)  # #b8dcc8
COR_VERDE_BG = RGBColor(0xDB, 0xEE, 0xE4)  # #dbeee4
COR_BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
COR_PRETO = RGBColor(0x0F, 0x17, 0x2A)
COR_CINZA_ESCURO = RGBColor(0x33, 0x41, 0x55)
COR_CINZA = RGBColor(0x64, 0x74, 0x8B)
COR_CINZA_CLARO = RGBColor(0xCB, 0xD5, 0xE1)
COR_CINZA_BG = RGBColor(0xF8, 0xFA, 0xFC)
COR_VERMELHO = RGBColor(0xDC, 0x26, 0x26)
COR_VERMELHO_BG = RGBColor(0xFE, 0xF2, 0xF2)
COR_VERDE_OK = RGBColor(0x15, 0x80, 0x3D)
COR_AMARELO_BG = RGBColor(0xFE, 0xF3, 0xC7)


# =============================================================================
# Domínio (mesmo da análise principal)
# =============================================================================


@dataclass
class Cafe:
    nome: str
    faturamento_1t: Decimal
    pessoal_1t: Decimal
    fte_total: int
    atendente: int
    aux_cozinha: int
    barista: int
    gerente: int
    lider: int

    @property
    def pessoal_mensal(self) -> Decimal:
        return (self.pessoal_1t / Decimal("3")).quantize(Decimal("0.01"))

    @property
    def perc_folha_fat(self) -> Decimal:
        if self.faturamento_1t == 0:
            return Decimal("0")
        return (self.pessoal_1t / self.faturamento_1t * 100).quantize(Decimal("0.01"))


def carregar_cafes() -> list[Cafe]:
    df = pd.read_excel(EXCEL_PATH, sheet_name="RESUMO CAFÉ", header=None)
    cafes: list[Cafe] = []
    for i in range(2, 12):
        row = df.iloc[i]
        nome = row[0]
        if pd.isna(nome) or "TOTAL" in str(nome).upper():
            continue
        cafes.append(
            Cafe(
                nome=str(nome).strip(),
                faturamento_1t=Decimal(str(row[10])) if pd.notna(row[10]) else Decimal("0"),
                pessoal_1t=Decimal(str(row[11])) if pd.notna(row[11]) else Decimal("0"),
                atendente=int(row[14]) if pd.notna(row[14]) else 0,
                aux_cozinha=int(row[15]) if pd.notna(row[15]) else 0,
                barista=int(row[16]) if pd.notna(row[16]) else 0,
                gerente=int(row[17]) if pd.notna(row[17]) else 0,
                lider=int(row[18]) if pd.notna(row[18]) else 0,
                fte_total=int(row[19]) if pd.notna(row[19]) else 0,
            )
        )
    return cafes


def simular(cafe: Cafe, cenario: str) -> dict:
    r = _ratio(cenario)
    fte_atual = Decimal(cafe.fte_total)
    fte_prop = Decimal(str(math.ceil(float(fte_atual * r) * 100) / 100))
    folha_atual = cafe.pessoal_mensal
    folha_prop = (folha_atual * r).quantize(Decimal("0.01"))
    delta = folha_prop - folha_atual
    return {
        "fte_proposto": fte_prop,
        "folha_prop": folha_prop,
        "delta_mes": delta,
        "delta_pct": (delta / folha_atual * 100) if folha_atual > 0 else Decimal("0"),
    }


# =============================================================================
# Formatação
# =============================================================================


def brl(v) -> str:  # type: ignore[no-untyped-def]
    return f"R$ {float(v):,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")


def brl_short(v) -> str:  # type: ignore[no-untyped-def]
    f = float(v)
    if abs(f) >= 1_000_000:
        return f"R$ {f / 1_000_000:.2f} mi".replace(".", ",")
    if abs(f) >= 1_000:
        return f"R$ {f / 1_000:.0f} mil"
    return f"R$ {f:.0f}"


def pct(v, decimals=2) -> str:  # type: ignore[no-untyped-def]
    return f"{float(v):.{decimals}f}%".replace(".", ",")


# =============================================================================
# Helpers de slide (slide widescreen 16:9 = 13.333 × 7.5 inches)
# =============================================================================


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color or fill_color
    shape.line.width = Pt(0)
    shape.shadow.inherit = False
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font_size=14,
    bold=False,
    color=COR_PRETO,
    align=PP_ALIGN.LEFT,
    font_name="Calibri",
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True

    # Suporte a múltiplas linhas
    lines = text.split("\n") if isinstance(text, str) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return tb


def add_brand_header(slide, title=None, subtitle=None):
    """Header verde fixo em todos os slides com logo e nome."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.5), COR_VERDE_PRIMARIO)
    add_text(
        slide,
        Inches(0.5),
        Inches(0.12),
        Inches(6),
        Inches(0.3),
        "MudAção Escala",
        font_size=12,
        bold=True,
        color=COR_VERDE_CLARO,
    )
    add_text(
        slide,
        Inches(7),
        Inches(0.12),
        Inches(6),
        Inches(0.3),
        "Análise PEC 8/2025 — Cafés T&F",
        font_size=11,
        color=COR_VERDE_CLARO,
        align=PP_ALIGN.RIGHT,
    )


def add_footer(slide, page_num, total_pages):
    """Footer com paginação + linha sutil."""
    add_rect(
        slide,
        Inches(0.5),
        Inches(7.05),
        SLIDE_W - Inches(1.0),
        Emu(6000),
        COR_CINZA_CLARO,
    )
    add_text(
        slide,
        Inches(0.5),
        Inches(7.15),
        Inches(6),
        Inches(0.3),
        "Confidencial · MudAção Escala · Maio/2026",
        font_size=9,
        color=COR_CINZA,
    )
    add_text(
        slide,
        Inches(7),
        Inches(7.15),
        Inches(6),
        Inches(0.3),
        f"{page_num} / {total_pages}",
        font_size=9,
        color=COR_CINZA,
        align=PP_ALIGN.RIGHT,
    )


# =============================================================================
# Slides
# =============================================================================


def slide_capa(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # BG verde escuro full
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, COR_VERDE_ESCURO)

    # Logo box (simples — um quadradinho com M)
    logo_left = Inches(0.7)
    logo_top = Inches(0.7)
    add_rect(slide, logo_left, logo_top, Inches(0.6), Inches(0.6), COR_VERDE_CLARO)
    add_text(
        slide,
        logo_left,
        logo_top + Inches(0.05),
        Inches(0.6),
        Inches(0.5),
        "M",
        font_size=28,
        bold=True,
        color=COR_VERDE_PRIMARIO,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        Inches(1.45),
        Inches(0.85),
        Inches(4),
        Inches(0.4),
        "MudAção",
        font_size=18,
        bold=True,
        color=COR_VERDE_CLARO,
    )

    # Pretitle
    add_text(
        slide,
        Inches(0.7),
        Inches(2.6),
        Inches(11),
        Inches(0.4),
        "ANÁLISE EXECUTIVA · 1º TRIMESTRE 2026",
        font_size=14,
        bold=True,
        color=COR_VERDE_CLARO,
    )

    # Título grande
    add_text(
        slide,
        Inches(0.7),
        Inches(3.1),
        Inches(11.5),
        Inches(1.2),
        "Impacto da PEC 8/2025\nna rede de cafés Track & Field",
        font_size=42,
        bold=True,
        color=COR_BRANCO,
    )

    # Subtitle
    add_text(
        slide,
        Inches(0.7),
        Inches(5.0),
        Inches(11.5),
        Inches(0.8),
        "Quanto custa a transição da escala 6x1 para 5x2 — e quanto dá pra recuperar",
        font_size=18,
        color=COR_VERDE_CLARO,
    )

    # Footer
    add_text(
        slide,
        Inches(0.7),
        Inches(6.7),
        Inches(11.5),
        Inches(0.3),
        "Preparado por MudAção · felipe@feldens.com · simulaescala.mudacao.com.br",
        font_size=12,
        color=COR_VERDE_CLARO,
    )


def slide_titulo(prs, page, total, kicker, titulo, subtitle=None):
    """Slide branco simples com título de seção."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_header(slide)

    add_text(
        slide,
        Inches(0.7),
        Inches(1.4),
        Inches(11),
        Inches(0.4),
        kicker,
        font_size=14,
        bold=True,
        color=COR_VERDE_PRIMARIO,
    )
    add_text(
        slide,
        Inches(0.7),
        Inches(1.85),
        Inches(11),
        Inches(0.8),
        titulo,
        font_size=32,
        bold=True,
        color=COR_VERDE_ESCURO,
    )
    if subtitle:
        add_text(
            slide,
            Inches(0.7),
            Inches(2.7),
            Inches(11),
            Inches(0.6),
            subtitle,
            font_size=16,
            color=COR_CINZA_ESCURO,
        )
    add_footer(slide, page, total)
    return slide


def slide_sumario(prs, page, total):
    slide = slide_titulo(prs, page, total, "SUMÁRIO", "O que você vai ver")
    items = [
        ("1", "Contexto: a PEC 8/2025"),
        ("2", "Estado atual da rede de cafés"),
        ("3", "Impacto financeiro projetado (3 cenários)"),
        ("4", "Cafés em risco estrutural"),
        ("5", "Quadro por função"),
        ("6", "Comparativo: cafés vs. lojas próprias"),
        ("7", "Recomendações específicas"),
        ("8", "Conclusões e próximos passos"),
    ]
    top = Inches(3.6)
    for i, (n, label) in enumerate(items):
        y = top + Inches(0.45 * i)
        add_text(slide, Inches(0.9), y, Inches(0.6), Inches(0.4), n,
                 font_size=20, bold=True, color=COR_VERDE_PRIMARIO)
        add_text(slide, Inches(1.5), y + Inches(0.05), Inches(10), Inches(0.4),
                 label, font_size=16, color=COR_CINZA_ESCURO)


def slide_contexto(prs, page, total):
    slide = slide_titulo(
        prs, page, total,
        "1 · CONTEXTO",
        "A PEC 8/2025 reestrutura a jornada de trabalho no Brasil",
    )
    # 3 cards explicando
    cards = [
        ("44h → 40h", "Redução de 4 horas semanais\n(jornada legal)"),
        ("6x1 → 5x2", "Eliminação da escala\n6 dias úteis + 1 folga"),
        ("8% a 14%", "Aumento típico de folha\nno varejo (Fitch)"),
    ]
    card_w = Inches(3.6)
    card_h = Inches(1.8)
    top = Inches(3.7)
    gap = Inches(0.4)
    total_w = card_w * 3 + gap * 2
    start_left = (SLIDE_W - total_w) / 2

    for i, (big, sub) in enumerate(cards):
        left = start_left + (card_w + gap) * i
        add_rect(slide, left, top, card_w, card_h, COR_VERDE_BG)
        add_rect(slide, left, top, card_w, Emu(50000), COR_VERDE_PRIMARIO)
        add_text(slide, left, top + Inches(0.3), card_w, Inches(0.7),
                 big, font_size=32, bold=True, color=COR_VERDE_PRIMARIO,
                 align=PP_ALIGN.CENTER)
        add_text(slide, left, top + Inches(1.05), card_w, Inches(0.8),
                 sub, font_size=13, color=COR_CINZA_ESCURO,
                 align=PP_ALIGN.CENTER)

    # Texto abaixo
    add_text(
        slide,
        Inches(1),
        Inches(5.9),
        Inches(11.3),
        Inches(0.6),
        "Para varejo e food service, cada estabelecimento precisa de mais pessoas pra cobrir "
        "o mesmo período de operação. A boa notícia: parte do impacto pode ser recuperado com "
        "Workforce Management baseado em IA (4–7%).",
        font_size=14,
        color=COR_CINZA_ESCURO,
    )


def slide_estado_atual(prs, page, total, cafes, totais):
    slide = slide_titulo(
        prs, page, total,
        "2 · ESTADO ATUAL",
        "A rede em números (1T26)",
        "10 cafés Track & Field analisados",
    )

    # 4 KPI cards no topo
    kpis = [
        ("10", "cafés", COR_VERDE_PRIMARIO),
        (brl_short(totais["fat"]), "faturamento 1T26", COR_VERDE_PRIMARIO),
        (brl_short(totais["folha_mes"]), "folha mensal média", COR_VERDE_PRIMARIO),
        (str(totais["fte"]), "FTEs (escala 6x1)", COR_VERDE_PRIMARIO),
    ]
    card_w = Inches(2.8)
    card_h = Inches(1.5)
    top = Inches(3.6)
    gap = Inches(0.2)
    total_w = card_w * 4 + gap * 3
    start = (SLIDE_W - total_w) / 2

    for i, (big, sub, c) in enumerate(kpis):
        left = start + (card_w + gap) * i
        add_rect(slide, left, top, card_w, card_h, COR_VERDE_BG)
        add_text(slide, left, top + Inches(0.3), card_w, Inches(0.6),
                 big, font_size=28, bold=True, color=c, align=PP_ALIGN.CENTER)
        add_text(slide, left, top + Inches(0.95), card_w, Inches(0.4),
                 sub, font_size=12, color=COR_CINZA_ESCURO, align=PP_ALIGN.CENTER)

    # Highlight box vermelho com alerta
    alert_top = Inches(5.4)
    add_rect(slide, Inches(1.5), alert_top, Inches(10.3), Inches(1.3), COR_VERMELHO_BG)
    add_rect(slide, Inches(1.5), alert_top, Inches(0.1), Inches(1.3), COR_VERMELHO)
    add_text(
        slide,
        Inches(1.8),
        alert_top + Inches(0.2),
        Inches(10),
        Inches(0.5),
        "⚠️  ALERTA  —  % FOLHA / FATURAMENTO = 56,12%",
        font_size=18,
        bold=True,
        color=COR_VERMELHO,
    )
    add_text(
        slide,
        Inches(1.8),
        alert_top + Inches(0.75),
        Inches(10),
        Inches(0.6),
        "Lojas próprias operam com 10,95%. Cafés estão com a margem 5× mais "
        "pressionada — qualquer aumento de folha bate direto no resultado.",
        font_size=13,
        color=COR_CINZA_ESCURO,
    )


def slide_impacto_3cenarios(prs, page, total, cafes, cenarios):
    slide = slide_titulo(
        prs, page, total,
        "3 · IMPACTO PROJETADO",
        "Três cenários de transição",
        "Variação do fator de produtividade aplicado sobre o ratio 44h/40h",
    )

    # Tabela com 3 colunas (cenários)
    cenario_cards = [
        ("PESSIMISTA", "Sem ganho de\nprodutividade",
         cenarios["pessimista"], COR_VERMELHO),
        ("NEUTRO", "Premissa padrão\n(estudo Fitch)",
         cenarios["neutro"], COR_VERDE_PRIMARIO),
        ("OTIMISTA", "Com WFM bem\nimplementado",
         cenarios["otimista"], COR_VERDE_OK),
    ]

    card_w = Inches(3.7)
    card_h = Inches(3.3)
    top = Inches(3.5)
    gap = Inches(0.3)
    total_w = card_w * 3 + gap * 2
    start = (SLIDE_W - total_w) / 2

    for i, (nome, desc, data, cor) in enumerate(cenario_cards):
        left = start + (card_w + gap) * i
        is_neutro = nome == "NEUTRO"

        bg_color = cor if is_neutro else COR_CINZA_BG
        text_color = COR_BRANCO if is_neutro else COR_CINZA_ESCURO
        accent_color = COR_BRANCO if is_neutro else cor

        add_rect(slide, left, top, card_w, card_h, bg_color)
        add_rect(slide, left, top, card_w, Inches(0.6), cor)

        # Header do card
        add_text(slide, left, top + Inches(0.15), card_w, Inches(0.4),
                 nome, font_size=14, bold=True, color=COR_BRANCO,
                 align=PP_ALIGN.CENTER)

        # Descrição
        add_text(slide, left, top + Inches(0.75), card_w, Inches(0.6),
                 desc, font_size=11, color=text_color, align=PP_ALIGN.CENTER)

        # Δ % grande
        add_text(slide, left, top + Inches(1.5), card_w, Inches(0.7),
                 f"+{pct(data['delta_pct'])}", font_size=32, bold=True,
                 color=accent_color, align=PP_ALIGN.CENTER)
        add_text(slide, left, top + Inches(2.15), card_w, Inches(0.3),
                 "Δ folha mensal", font_size=11, color=text_color,
                 align=PP_ALIGN.CENTER)

        # Anual
        add_text(slide, left, top + Inches(2.55), card_w, Inches(0.5),
                 brl_short(data["delta_mes"] * 12), font_size=18, bold=True,
                 color=accent_color, align=PP_ALIGN.CENTER)
        add_text(slide, left, top + Inches(3.0), card_w, Inches(0.3),
                 "no ano", font_size=11, color=text_color,
                 align=PP_ALIGN.CENTER)


def slide_cenario_neutro(prs, page, total, cafes, neutro):
    slide = slide_titulo(
        prs, page, total,
        "4 · CENÁRIO NEUTRO",
        "Projeção pós-PEC 8 (premissa Fitch)",
        "O que esperar se nada mudar na operação",
    )

    # Headline grande
    add_text(
        slide,
        Inches(0.7),
        Inches(3.6),
        Inches(12),
        Inches(0.8),
        f"{brl_short(neutro['delta_mes'])}  / mês",
        font_size=44,
        bold=True,
        color=COR_VERDE_PRIMARIO,
    )
    add_text(
        slide,
        Inches(0.7),
        Inches(4.5),
        Inches(12),
        Inches(0.5),
        "de aumento de folha na rede de cafés (cenário neutro)",
        font_size=16,
        color=COR_CINZA_ESCURO,
    )

    # 3 colunas embaixo
    items = [
        (brl_short(neutro["delta_mes"] * 12), "no ano"),
        (f"+{float(neutro['fte_extras']):.1f}", "FTEs extras pra contratar"),
        (brl_short(neutro["economia_wfm_ano"]), "recuperáveis com WFM (ano)"),
    ]
    card_w = Inches(3.6)
    top = Inches(5.5)
    gap = Inches(0.4)
    total_w = card_w * 3 + gap * 2
    start = (SLIDE_W - total_w) / 2

    for i, (big, sub) in enumerate(items):
        left = start + (card_w + gap) * i
        add_rect(slide, left, top, card_w, Inches(1.2), COR_VERDE_BG)
        add_text(slide, left, top + Inches(0.2), card_w, Inches(0.5),
                 big, font_size=22, bold=True, color=COR_VERDE_PRIMARIO,
                 align=PP_ALIGN.CENTER)
        add_text(slide, left, top + Inches(0.8), card_w, Inches(0.4),
                 sub, font_size=11, color=COR_CINZA_ESCURO,
                 align=PP_ALIGN.CENTER)


def slide_risco(prs, page, total, em_risco):
    slide = slide_titulo(
        prs, page, total,
        "5 · CAFÉS EM RISCO ESTRUTURAL",
        f"{len(em_risco)} cafés operam com folha > 70% do faturamento",
        "Pra estes, o impacto da PEC 8 é potencialmente terminal",
    )

    # Tabela manual
    top = Inches(3.7)
    row_h = Inches(0.55)
    col_widths = [Inches(4.2), Inches(2.2), Inches(2.2), Inches(2.4)]
    headers = ["Café", "Fat 1T26", "Folha/mês", "% folha/fat"]
    start_left = Inches(1.1)

    # Header
    x = start_left
    add_rect(slide, x, top, sum(col_widths, Inches(0)), row_h, COR_VERDE_PRIMARIO)
    for i, (h, w) in enumerate(zip(headers, col_widths)):
        add_text(slide, x + Inches(0.15), top + Inches(0.13), w, Inches(0.4),
                 h, font_size=12, bold=True, color=COR_BRANCO,
                 align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT)
        x += w

    # Rows
    for ridx, c in enumerate(em_risco):
        row_top = top + row_h + row_h * ridx
        bg = COR_BRANCO if ridx % 2 == 0 else COR_CINZA_BG
        add_rect(slide, start_left, row_top, sum(col_widths, Inches(0)), row_h, bg)

        cells = [
            c.nome,
            brl_short(c.faturamento_1t),
            brl_short(c.pessoal_mensal),
            pct(c.perc_folha_fat),
        ]
        x = start_left
        for i, (val, w) in enumerate(zip(cells, col_widths)):
            color = COR_VERMELHO if i == 3 else COR_CINZA_ESCURO
            bold = i == 0 or i == 3
            add_text(
                slide, x + Inches(0.15), row_top + Inches(0.15),
                w - Inches(0.3), Inches(0.4),
                val, font_size=12, bold=bold, color=color,
                align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT,
            )
            x += w

    # Insight
    insight_top = top + row_h * (len(em_risco) + 1) + Inches(0.3)
    add_rect(slide, Inches(1.1), insight_top, Inches(11.1), Inches(1), COR_AMARELO_BG)
    add_text(
        slide, Inches(1.3), insight_top + Inches(0.2), Inches(10.7), Inches(0.6),
        "TFC Rio Sul gasta R$ 1,11 em folha pra cada R$ 1,00 que fatura — está com prejuízo "
        "operacional antes mesmo da PEC. Pós-PEC fica inviável sem intervenção.",
        font_size=12, color=COR_CINZA_ESCURO,
    )


def slide_quadro(prs, page, total, cafes):
    slide = slide_titulo(
        prs, page, total,
        "6 · QUADRO POR FUNÇÃO",
        "Composição do headcount na rede (94 FTEs)",
    )
    totals = {
        "Barista": sum(c.barista for c in cafes),
        "Atendente": sum(c.atendente for c in cafes),
        "Gerente de Loja": sum(c.gerente for c in cafes),
        "Líder": sum(c.lider for c in cafes),
        "Auxiliar de Cozinha": sum(c.aux_cozinha for c in cafes),
    }
    total_fte = sum(totals.values())

    top = Inches(3.6)
    bar_max_w = Inches(7)
    for i, (func, val) in enumerate(totals.items()):
        y = top + Inches(0.65 * i)
        pct_val = val / total_fte * 100

        # Nome
        add_text(slide, Inches(1.1), y, Inches(2.5), Inches(0.4),
                 func, font_size=14, color=COR_CINZA_ESCURO, bold=True)

        # Barra background
        add_rect(slide, Inches(3.7), y + Inches(0.05), bar_max_w, Inches(0.3),
                 COR_CINZA_BG)

        # Barra colorida (proporcional ao % do total)
        bar_w = Emu(int(bar_max_w * (val / max(totals.values()))))
        add_rect(slide, Inches(3.7), y + Inches(0.05), bar_w, Inches(0.3),
                 COR_VERDE_PRIMARIO)

        # Valor + %
        add_text(slide, Inches(11), y, Inches(1.5), Inches(0.4),
                 f"{val} ({pct_val:.0f}%)", font_size=14, bold=True,
                 color=COR_VERDE_PRIMARIO)

    # Insight
    insight_top = Inches(7.0) - Inches(1.6)
    add_rect(slide, Inches(1.1), insight_top, Inches(11.1), Inches(1), COR_VERDE_BG)
    add_text(
        slide, Inches(1.3), insight_top + Inches(0.2), Inches(10.7), Inches(0.6),
        "75% do quadro é Barista + Atendente — funções com forte overlap operacional. "
        "Multifuncionalidade pode reduzir contratações pós-PEC sem aumentar carga.",
        font_size=12, color=COR_CINZA_ESCURO,
    )


def slide_comparativo(prs, page, total, totais, neutro):
    slide = slide_titulo(
        prs, page, total,
        "7 · COMPARATIVO",
        "Cafés vs. Lojas Próprias (1T26)",
        "Mesma metodologia, leituras muito diferentes",
    )

    # Tabela manual com 3 colunas
    top = Inches(3.5)
    row_h = Inches(0.45)
    col_widths = [Inches(4.7), Inches(3.0), Inches(3.0)]

    rows = [
        ("Métrica", "CAFÉS", "LOJAS PRÓPRIAS", True),
        ("Unidades", "10", "55", False),
        ("Faturamento 1T26", brl_short(totais["fat"]), "R$ 192,34 mi", False),
        ("Folha 1T26", brl_short(totais["folha_1t"]), "R$ 21,07 mi", False),
        ("% folha / faturamento", pct(totais["pct_folha"]), "10,95%", False),
        ("FTEs atual (6x1)", str(totais["fte"]), "728", False),
        ("Δ folha mensal (neutro)", brl_short(neutro["delta_mes"]), "R$ 683 mil", False),
        ("Δ folha anual (neutro)", brl_short(neutro["delta_mes"] * 12), "R$ 8,20 mi", False),
        ("Δ anual / Faturamento anual", pct(neutro["impacto_fat_pct"]), "1,07%", False),
    ]

    start_left = Inches(1.1)
    for ridx, (label, val_c, val_l, is_header) in enumerate(rows):
        row_top = top + row_h * ridx
        if is_header:
            add_rect(slide, start_left, row_top, sum(col_widths, Inches(0)), row_h,
                     COR_VERDE_PRIMARIO)
            color = COR_BRANCO
        else:
            bg = COR_BRANCO if ridx % 2 == 0 else COR_CINZA_BG
            add_rect(slide, start_left, row_top, sum(col_widths, Inches(0)), row_h, bg)
            color = COR_CINZA_ESCURO

        x = start_left
        cells = [label, val_c, val_l]
        for i, (val, w) in enumerate(zip(cells, col_widths)):
            cell_color = color
            # Destaque vermelho na linha % folha/fat coluna cafés
            if not is_header and label == "% folha / faturamento" and i == 1:
                cell_color = COR_VERMELHO
            if not is_header and "Δ anual" in label and i == 1:
                cell_color = COR_VERMELHO

            add_text(
                slide, x + Inches(0.15), row_top + Inches(0.1),
                w - Inches(0.3), Inches(0.4),
                val, font_size=12,
                bold=is_header or i == 0,
                color=cell_color,
                align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT,
            )
            x += w


def slide_insight(prs, page, total, neutro):
    slide = slide_titulo(
        prs, page, total,
        "8 · INSIGHT CENTRAL",
        "Onde dói mais não é o valor absoluto",
    )

    # Caixa grande destacada
    box_top = Inches(3.5)
    box_h = Inches(3.0)
    add_rect(slide, Inches(1.1), box_top, Inches(11.1), box_h, COR_VERDE_ESCURO)
    add_rect(slide, Inches(1.1), box_top, Inches(0.15), box_h, COR_VERDE_CLARO)

    add_text(
        slide,
        Inches(1.5),
        box_top + Inches(0.4),
        Inches(10.5),
        Inches(0.5),
        "EM VALOR ABSOLUTO,",
        font_size=12, bold=True, color=COR_VERDE_CLARO,
    )
    add_text(
        slide,
        Inches(1.5),
        box_top + Inches(0.85),
        Inches(10.5),
        Inches(0.7),
        "lojas próprias sofrem ~13× mais (R$ 8,2 mi vs R$ 607 mil/ano).",
        font_size=22, color=COR_BRANCO,
    )

    add_text(
        slide,
        Inches(1.5),
        box_top + Inches(1.7),
        Inches(10.5),
        Inches(0.5),
        "EM % SOBRE O FATURAMENTO,",
        font_size=12, bold=True, color=COR_VERDE_CLARO,
    )
    add_text(
        slide,
        Inches(1.5),
        box_top + Inches(2.15),
        Inches(10.5),
        Inches(0.7),
        "cafés sofrem 5× mais.",
        font_size=26, bold=True, color=COR_BRANCO,
    )


def slide_recomendacoes(prs, page, total):
    slide = slide_titulo(
        prs, page, total,
        "9 · RECOMENDAÇÕES ESPECÍFICAS",
        "Onde os cafés precisam atuar (e não basta absorver)",
    )

    recs = [
        ("1", "Revisar viabilidade dos 3 cafés críticos",
         "Rio Sul, Catarina, Vila Lobos — antes da PEC entrar em vigor"),
        ("2", "Reduzir headcount via multifuncionalidade",
         "Barista + Atendente = 75% do quadro, há overlap natural"),
        ("3", "Repensar horário de operação",
         "Café tem curva de demanda mais nítida que varejo — fechar em pico baixo"),
        ("4", "Modelo híbrido CLT + horista nas pontas",
         "Reduz CLT fixo, paga só horas necessárias em fim-de-semana"),
        ("5", "Implementar WFM com IA",
         f"5% de economia da folha proposta = R$ 343 mil/ano recuperados"),
        ("6", "Renegociação coletiva específica",
         "Sindicato de gastronomia pode oferecer flexibilidade ao setor pressionado"),
    ]

    top = Inches(3.4)
    cols = 2
    rows_count = 3
    card_w = Inches(5.6)
    card_h = Inches(1.1)
    h_gap = Inches(0.3)
    v_gap = Inches(0.25)
    total_w = card_w * cols + h_gap * (cols - 1)
    start_left = (SLIDE_W - total_w) / 2

    for idx, (n, titulo, desc) in enumerate(recs):
        col = idx % cols
        row = idx // cols
        left = start_left + (card_w + h_gap) * col
        y = top + (card_h + v_gap) * row

        add_rect(slide, left, y, card_w, card_h, COR_CINZA_BG)
        add_rect(slide, left, y, Inches(0.1), card_h, COR_VERDE_PRIMARIO)
        # Número grande
        add_text(slide, left + Inches(0.2), y + Inches(0.15), Inches(0.5), Inches(0.6),
                 n, font_size=28, bold=True, color=COR_VERDE_PRIMARIO)
        # Título
        add_text(slide, left + Inches(0.85), y + Inches(0.18), card_w - Inches(1), Inches(0.4),
                 titulo, font_size=13, bold=True, color=COR_VERDE_ESCURO)
        # Descrição
        add_text(slide, left + Inches(0.85), y + Inches(0.6), card_w - Inches(1), Inches(0.5),
                 desc, font_size=10, color=COR_CINZA_ESCURO)


def slide_conclusao(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, COR_VERDE_ESCURO)

    add_text(
        slide,
        Inches(0.7),
        Inches(0.7),
        Inches(11),
        Inches(0.4),
        "CONCLUSÕES E PRÓXIMOS PASSOS",
        font_size=14, bold=True, color=COR_VERDE_CLARO,
    )

    add_text(
        slide,
        Inches(0.7),
        Inches(1.4),
        Inches(11.5),
        Inches(1.0),
        "Não há ‘absorver’ a PEC 8 nos cafés.",
        font_size=38, bold=True, color=COR_BRANCO,
    )
    add_text(
        slide,
        Inches(0.7),
        Inches(2.45),
        Inches(11.5),
        Inches(0.7),
        "É um momento de redesenhar a operação.",
        font_size=22, color=COR_VERDE_CLARO,
    )

    bullets = [
        "Os 3 cafés em risco estrutural exigem decisão imediata (antes da PEC)",
        "WFM com IA pode recuperar R$ 343 mil/ano — vale o investimento",
        "Multifuncionalidade Barista ↔ Atendente reduz contratações de 9 FTEs",
        "Renegociação coletiva específica pode dar flexibilidade adicional",
    ]
    top = Inches(3.7)
    for i, b in enumerate(bullets):
        y = top + Inches(0.55 * i)
        add_text(slide, Inches(0.95), y + Inches(0.05), Inches(0.3), Inches(0.4),
                 "▸", font_size=20, bold=True, color=COR_VERDE_CLARO)
        add_text(slide, Inches(1.3), y, Inches(11), Inches(0.5),
                 b, font_size=15, color=COR_BRANCO)

    # Footer com contato
    add_rect(slide, 0, Inches(6.5), SLIDE_W, Inches(1), COR_VERDE_PRIMARIO)
    add_text(
        slide,
        Inches(0.7),
        Inches(6.75),
        Inches(11.5),
        Inches(0.4),
        "Análise preparada por MudAção · Felipe Feldens",
        font_size=14, bold=True, color=COR_BRANCO,
    )
    add_text(
        slide,
        Inches(0.7),
        Inches(7.1),
        Inches(11.5),
        Inches(0.3),
        "felipe@feldens.com  ·  simulaescala.mudacao.com.br  ·  (11) 99632-5174",
        font_size=12, color=COR_VERDE_CLARO,
    )


# =============================================================================
# Build
# =============================================================================


def gerar() -> Path:
    cafes = carregar_cafes()

    # Roda os 3 cenários
    cenarios = {}
    for cen in ("pessimista", "neutro", "otimista"):
        sims = [simular(c, cen) for c in cafes]
        cenarios[cen] = {
            "fte_total": sum(s["fte_proposto"] for s in sims),
            "folha_total": sum(s["folha_prop"] for s in sims),
            "delta_mes": sum(s["delta_mes"] for s in sims),
            "delta_pct": (
                sum(s["delta_mes"] for s in sims)
                / sum(c.pessoal_mensal for c in cafes)
                * 100
            ),
        }

    # Totais
    fat_total = sum(c.faturamento_1t for c in cafes)
    folha_1t = sum(c.pessoal_1t for c in cafes)
    folha_mes = folha_1t / Decimal("3")
    fte_atual = sum(c.fte_total for c in cafes)
    totais = {
        "fat": fat_total,
        "folha_1t": folha_1t,
        "folha_mes": folha_mes,
        "pct_folha": (folha_1t / fat_total * 100) if fat_total > 0 else Decimal("0"),
        "fte": fte_atual,
    }

    # Neutro detalhado
    neutro = cenarios["neutro"]
    neutro["fte_extras"] = neutro["fte_total"] - Decimal(fte_atual)
    neutro["economia_wfm_ano"] = neutro["folha_total"] * WFM_ECONOMY_PCT * 12
    neutro["impacto_fat_pct"] = (
        neutro["delta_mes"] * 12 / (fat_total * 4) * 100
    ) if fat_total > 0 else Decimal("0")

    em_risco = sorted(
        [c for c in cafes if c.perc_folha_fat > 70],
        key=lambda c: c.perc_folha_fat,
        reverse=True,
    )

    # =========================================================================
    # Build PPTX
    # =========================================================================
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total_slides = 11

    slide_capa(prs)
    slide_sumario(prs, 2, total_slides)
    slide_contexto(prs, 3, total_slides)
    slide_estado_atual(prs, 4, total_slides, cafes, totais)
    slide_impacto_3cenarios(prs, 5, total_slides, cafes, cenarios)
    slide_cenario_neutro(prs, 6, total_slides, cafes, neutro)
    slide_risco(prs, 7, total_slides, em_risco)
    slide_quadro(prs, 8, total_slides, cafes)
    slide_comparativo(prs, 9, total_slides, totais, neutro)
    slide_insight(prs, 10, total_slides, neutro)
    slide_recomendacoes(prs, 11, total_slides)
    slide_conclusao(prs, 12, total_slides)

    out_dir = Path(__file__).resolve().parents[1] / "tmp"
    out_dir.mkdir(exist_ok=True)
    out_path = out_dir / "Apresentacao_TF_Cafes_PEC8.pptx"
    prs.save(out_path)
    return out_path


def main():
    print("🎨 Gerando apresentação MudAção Escala — Cafés T&F...")
    path = gerar()
    print(f"\n✓ Apresentação salva em:\n  {path}")
    print(f"\n  Abre direto: start \"\" \"{path}\"")
    print(f"  Tamanho: {path.stat().st_size / 1024:.0f} KB")


if __name__ == "__main__":
    main()
